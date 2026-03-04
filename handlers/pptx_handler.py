"""Generate PowerPoint (.pptx) from HTML for rich, colorful, interactive slides."""

import base64
import re
from io import BytesIO

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from utils.security import validate_image_path
from utils.file_utils import generate_output_path


# Dimensões típicas do slide (polegadas)
SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5
MARGIN = 0.5
CONTENT_TOP = 1.2
CONTENT_WIDTH = SLIDE_WIDTH - 2 * MARGIN
TITLE_HEIGHT = 0.8


def create_pptx_from_html(
    html_content: str,
    output_path: str,
    title: str | None = None,
) -> str:
    """Generate a PowerPoint presentation from HTML.

    Supports:
      - Multiple slides: use <section> or <div class="slide"> for each slide, or split by <h1>
      - Rich text: <b>, <strong>, <i>, <em>, <span style="color: #hex or rgb()">
      - Tables: <table>, <thead>, <tbody>, <th>, <td>
      - Lists: <ul>, <ol>, <li>
      - Images: <img src="file.png"> or data URI (data:image/png;base64,...)
      - Links: <a href="url"> become clickable hyperlinks

    Returns the absolute path of the generated .pptx file.
    """
    soup = BeautifulSoup(html_content, "html.parser")
    slide_blocks = _extract_slide_blocks(soup)
    if not slide_blocks:
        slide_blocks = [soup]  # single slide with full body

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)
    blank_layout = prs.slide_layouts[6]  # blank

    for idx, block in enumerate(slide_blocks):
        slide = prs.slides.add_slide(blank_layout)
        use_presentation_title = title if (idx == 0 and title) else None
        _render_slide(slide, block, use_presentation_title)

    dest = generate_output_path(output_path)
    if not str(dest).lower().endswith(".pptx"):
        dest = dest.with_suffix(".pptx")
    prs.save(str(dest))
    return str(dest)


def _extract_slide_blocks(soup: BeautifulSoup) -> list:
    """Split HTML into slide blocks: <section> or <div class=\"slide\">. Otherwise one slide."""
    for selector in ("section", "div.slide", "[data-slide]"):
        elements = soup.select(selector)
        if elements:
            return elements
    return []


def _render_slide(slide, block, presentation_title: str | None):
    """Render one slide from a BeautifulSoup node."""
    container = block.find("body") or block if hasattr(block, "find") else block
    if not container:
        return
    # Title: use presentation_title on first slide if given, else first h1/h2 in block
    title_el = container.find(["h1", "h2"])
    title_text = presentation_title or (title_el.get_text(strip=True) if title_el else "")
    if title_text:
        _add_title_shape(slide, title_text)
        y = Inches(CONTENT_TOP + TITLE_HEIGHT)
    else:
        y = Inches(CONTENT_TOP)

    left = Inches(MARGIN)
    width = Inches(CONTENT_WIDTH)
    max_height = SLIDE_HEIGHT - (y.inches if hasattr(y, "inches") else CONTENT_TOP + TITLE_HEIGHT) - MARGIN

    for el in container.children:
        if not hasattr(el, "name") or el.name is None:
            continue
        if el.name in ("h1", "h2") and title_text:
            continue  # já usado como título
        y = _add_element(slide, el, left, y, width, max_height)
        if y is None:
            break


def _add_title_shape(slide, text: str):
    """Add title text box at top of slide."""
    title_box = slide.shapes.add_textbox(
        Inches(MARGIN),
        Inches(0.35),
        Inches(CONTENT_WIDTH),
        Inches(TITLE_HEIGHT),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)


def _add_element(slide, el, left, y_pos, width, max_height):
    """Add one HTML element to slide; return new y position (Inches) or None."""
    from pptx.util import Emu

    if el.name in ("h1", "h2", "h3", "h4"):
        level = int(el.name[1])
        size = Pt(24 - (level - 1) * 2)
        box = slide.shapes.add_textbox(left, y_pos, width, Inches(0.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _add_runs_from_node(p, el)
        p.font.size = size
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
        return y_pos + Inches(0.55)

    if el.name == "p":
        box = slide.shapes.add_textbox(left, y_pos, width, Inches(0.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _add_runs_from_node(p, el)
        p.font.size = Pt(14)
        return y_pos + Inches(0.4)

    if el.name == "table":
        rows, cols = _table_dim(el)
        if rows and cols:
            table_height = min(Inches(0.4 * rows), Inches(max_height * 0.6))
            table_shape = slide.shapes.add_table(
                rows, cols,
                left, y_pos,
                width, table_height,
            )
            _fill_table(table_shape.table, el)
            return y_pos + table_height + Inches(0.2)

    if el.name in ("ul", "ol"):
        items = el.find_all("li", recursive=False)
        box = slide.shapes.add_textbox(left, y_pos, width, Inches(0.4 * len(items) + 0.2))
        tf = box.text_frame
        tf.word_wrap = True
        for i, li in enumerate(items):
            p = tf.add_paragraph() if i else tf.paragraphs[0]
            bullet = "• " if el.name == "ul" else f"{i + 1}. "
            p.text = bullet + li.get_text(separator=" ", strip=True)
            p.font.size = Pt(12)
            p.space_after = Pt(4)
        return y_pos + Inches(0.4 * len(items) + 0.3)

    if el.name == "img":
        src = el.get("src") or ""
        try:
            img_path_or_stream = _resolve_image_src(src)
            if img_path_or_stream:
                pic = slide.shapes.add_picture(
                    img_path_or_stream,
                    left, y_pos,
                    width=min(width, Inches(4)),
                )
                height_inches = pic.height / 914400.0
                return Inches(y_pos.inches + height_inches + 0.15)
        except Exception:
            pass
        return y_pos

    # div: recurse into children
    if el.name == "div":
        for child in el.children:
            if hasattr(child, "name") and child.name:
                y_pos = _add_element(slide, child, left, y_pos, width, max_height)
                if y_pos is None:
                    return None
        return y_pos

    return y_pos


def _add_runs_from_node(paragraph, node):
    """Append runs to paragraph from HTML node, preserving bold/italic/color/links."""
    default_size = paragraph.font.size or Pt(14)
    for content in node.children:
        if isinstance(content, str):
            if content.strip():
                run = paragraph.add_run()
                run.text = content
                run.font.size = default_size
            continue
        if not hasattr(content, "name") or content.name is None:
            continue
        if content.name in ("b", "strong"):
            run = paragraph.add_run()
            run.text = content.get_text()
            run.font.bold = True
            run.font.size = default_size
        elif content.name in ("i", "em"):
            run = paragraph.add_run()
            run.text = content.get_text()
            run.font.italic = True
            run.font.size = default_size
        elif content.name == "a":
            run = paragraph.add_run()
            run.text = content.get_text()
            run.font.size = default_size
            href = content.get("href")
            if href and hasattr(run, "hyperlink"):
                run.hyperlink.address = href
            try:
                run.font.color.rgb = RGBColor(0x34, 0x98, 0xDB)
            except Exception:
                pass
        elif content.name == "span":
            run = paragraph.add_run()
            run.text = content.get_text()
            run.font.size = default_size
            color = _parse_style_color(content.get("style", ""))
            if color:
                try:
                    run.font.color.rgb = RGBColor(*color)
                except Exception:
                    pass
        else:
            _add_runs_from_node(paragraph, content)


def _parse_style_color(style: str) -> tuple[int, int, int] | None:
    """Parse color from style string; return (r, g, b) or None."""
    if not style:
        return None
    # color: #fff or #ffffff
    m = re.search(r"color\s*:\s*#([0-9a-fA-F]{3}|[0-9a-fA-F]{6})\b", style)
    if m:
        hex_val = m.group(1)
        if len(hex_val) == 3:
            hex_val = "".join(c * 2 for c in hex_val)
        return (int(hex_val[0:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16))
    # color: rgb( r, g, b )
    m = re.search(r"color\s*:\s*rgb\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)", style)
    if m:
        return (int(m.group(1)), int(m.group(2)), int(m.group(3)))
    return None


def _table_dim(table_el) -> tuple[int, int]:
    """Return (rows, cols) for a table element."""
    rows_el = table_el.find("tbody") or table_el
    rows = rows_el.find_all("tr")
    if not rows:
        return 0, 0
    cols = max(len(r.find_all(["th", "td"])) for r in rows)
    return len(rows), cols


def _fill_table(table, table_el):
    """Fill python-pptx table from HTML table."""
    thead = table_el.find("thead")
    tbody = table_el.find("tbody") or table_el
    row_els = (thead.find_all("tr") if thead else []) + tbody.find_all("tr")
    for i, tr in enumerate(row_els):
        cells = tr.find_all(["th", "td"])
        for j, cell in enumerate(cells):
            if i < len(table.rows) and j < len(table.columns):
                table.cell(i, j).text = cell.get_text(separator=" ", strip=True)
                if cell.name == "th":
                    try:
                        table.cell(i, j).text_frame.paragraphs[0].font.bold = True
                        table.cell(i, j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
                    except Exception:
                        pass


def _resolve_image_src(src: str):
    """Return file path or file-like object for img src (path or data URI)."""
    src = (src or "").strip()
    if not src:
        return None
    if src.startswith("data:"):
        # data:image/png;base64,...
        m = re.match(r"data:image/[^;]+;base64,(.+)", src)
        if m:
            try:
                data = base64.b64decode(m.group(1))
                return BytesIO(data)
            except Exception:
                return None
        return None
    # path: validate and use
    try:
        path = validate_image_path(src)
        return str(path)
    except Exception:
        return None
